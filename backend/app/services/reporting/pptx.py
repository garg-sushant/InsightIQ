"""PowerPoint deck export (python-pptx).

Built from blank layouts rather than the default template placeholders, so the
deck looks deliberate rather than like an unedited Office starter file.
"""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.schemas.ai import InsightBundle
from app.schemas.analytics import AnalyticsResult, RiskLevel
from app.services.reporting.markdown import plain_bullets

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

INK = RGBColor(0x1F, 0x29, 0x33)
MUTED = RGBColor(0x61, 0x6E, 0x7C)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF5, 0xF7, 0xFA)
POSITIVE = RGBColor(0x05, 0x96, 0x69)
NEGATIVE = RGBColor(0xDC, 0x26, 0x26)

_LEVEL_COLORS = {
    RiskLevel.OK: POSITIVE,
    RiskLevel.WATCH: RGBColor(0xD9, 0x77, 0x06),
    RiskLevel.ELEVATED: RGBColor(0xEA, 0x58, 0x0C),
    RiskLevel.CRITICAL: NEGATIVE,
}

MARGIN = Inches(0.6)
CONTENT_WIDTH = SLIDE_WIDTH - 2 * MARGIN


def _blank(presentation: Presentation):
    # Layout 6 is the blank layout in the default template.
    return presentation.slides.add_slide(presentation.slide_layouts[6])


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    return frame


def _set(paragraph, text: str, *, size: int, bold: bool = False,
         color: RGBColor = INK, align=PP_ALIGN.LEFT) -> None:
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _title_slide(presentation: Presentation, result: AnalyticsResult,
                 organization_name: str, title: str) -> None:
    slide = _blank(presentation)
    band = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_WIDTH, Inches(2.6))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    band.shadow.inherit = False

    frame = _textbox(slide, MARGIN, Inches(0.8), CONTENT_WIDTH, Inches(1.6))
    _set(frame.paragraphs[0], title, size=40, bold=True, color=WHITE)
    _set(frame.add_paragraph(), organization_name, size=18, color=WHITE)

    period = result.period
    details = _textbox(slide, MARGIN, Inches(3.2), CONTENT_WIDTH, Inches(2.4))
    _set(
        details.paragraphs[0],
        f"{period.start:%d %B %Y} – {period.end:%d %B %Y}  ·  {period.days} days",
        size=20, bold=True,
    )
    _set(
        details.add_paragraph(),
        f"{result.row_count:,} order lines analysed  ·  Business health "
        f"{result.health.score}/100 (grade {result.health.grade})",
        size=15, color=MUTED,
    )
    _set(
        details.add_paragraph(),
        "All figures computed deterministically by InsightIQ from source data.",
        size=12, color=MUTED,
    )


def _kpi_slide(presentation: Presentation, result: AnalyticsResult) -> None:
    slide = _blank(presentation)
    _heading(slide, "Performance at a glance")

    kpis = result.kpis
    cards = [
        ("Revenue", f"${kpis.revenue.current:,.0f}", kpis.revenue.delta_pct, "%", True),
        ("Profit", f"${kpis.profit.current:,.0f}", kpis.profit.delta_pct, "%", True),
        ("Margin", f"{kpis.margin_pct.current:,.1f}%", kpis.margin_pct.delta_abs, " pp", True),
        ("Orders", f"{kpis.orders.current:,.0f}", kpis.orders.delta_pct, "%", True),
        ("Avg order value", f"${kpis.aov.current:,.0f}", kpis.aov.delta_pct, "%", True),
        ("Units sold", f"{kpis.units.current:,.0f}", kpis.units.delta_pct, "%", True),
        ("Return rate", f"{kpis.return_rate_pct.current:,.1f}%",
         kpis.return_rate_pct.delta_abs, " pp", False),
        ("Repeat rate", f"{kpis.repeat_rate_pct.current:,.1f}%",
         kpis.repeat_rate_pct.delta_abs, " pp", True),
    ]

    columns, gap = 4, Inches(0.25)
    card_width = int((CONTENT_WIDTH - gap * (columns - 1)) / columns)
    card_height = Inches(1.7)

    for index, (label, value, delta, unit, higher_is_better) in enumerate(cards):
        row, column = divmod(index, columns)
        left = MARGIN + column * (card_width + gap)
        top = Inches(1.6) + row * (card_height + Inches(0.3))

        shape = slide.shapes.add_shape(5, left, top, card_width, card_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD
        shape.line.color.rgb = RGBColor(0xE4, 0xE7, 0xEB)
        shape.shadow.inherit = False

        frame = _textbox(slide, left + Inches(0.22), top + Inches(0.18),
                         card_width - Inches(0.44), card_height - Inches(0.36))
        _set(frame.paragraphs[0], label.upper(), size=10, color=MUTED)
        _set(frame.add_paragraph(), value, size=24, bold=True)

        if delta is None:
            _set(frame.add_paragraph(), "no comparison period", size=10, color=MUTED)
        else:
            improving = (delta > 0) is higher_is_better
            _set(
                frame.add_paragraph(),
                f"{delta:+.1f}{unit} vs previous",
                size=11,
                bold=True,
                color=POSITIVE if improving else NEGATIVE,
            )


def _heading(slide, text: str, subtitle: str | None = None) -> None:
    frame = _textbox(slide, MARGIN, Inches(0.45), CONTENT_WIDTH, Inches(1.0))
    _set(frame.paragraphs[0], text, size=26, bold=True)
    if subtitle:
        _set(frame.add_paragraph(), subtitle, size=12, color=MUTED)


def _chart_slide(presentation: Presentation, title: str, images: list[bytes],
                 subtitle: str | None = None) -> None:
    slide = _blank(presentation)
    _heading(slide, title, subtitle)

    top = Inches(1.7)
    available_height = SLIDE_HEIGHT - top - Inches(0.5)

    if len(images) == 1:
        slide.shapes.add_picture(io.BytesIO(images[0]), MARGIN, top,
                                 width=CONTENT_WIDTH)
    else:
        gap = Inches(0.3)
        width = int((CONTENT_WIDTH - gap) / 2)
        for index, image in enumerate(images[:2]):
            left = MARGIN + index * (width + gap)
            slide.shapes.add_picture(io.BytesIO(image), left, top, width=width)
    del available_height


def _bullet_slide(presentation: Presentation, title: str, bullets: list[str],
                  subtitle: str | None = None) -> None:
    slide = _blank(presentation)
    _heading(slide, title, subtitle)

    frame = _textbox(slide, MARGIN, Inches(1.8), CONTENT_WIDTH,
                     SLIDE_HEIGHT - Inches(2.4))
    if not bullets:
        _set(frame.paragraphs[0], "No content available for this section.",
             size=14, color=MUTED)
        return

    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        _set(paragraph, f"•  {bullet}", size=14)
        paragraph.space_after = Pt(10)


def _health_slide(presentation: Presentation, result: AnalyticsResult) -> None:
    slide = _blank(presentation)
    health = result.health
    _heading(
        slide,
        f"Business health: {health.score}/100 (grade {health.grade})",
        health.headline,
    )

    frame = _textbox(slide, MARGIN, Inches(1.9), CONTENT_WIDTH, Inches(5.0))
    for index, indicator in enumerate(health.indicators):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = ""
        label_run = paragraph.add_run()
        label_run.text = f"{indicator.level.value.upper():<9} "
        label_run.font.size = Pt(12)
        label_run.font.bold = True
        label_run.font.name = "Consolas"
        label_run.font.color.rgb = _LEVEL_COLORS[indicator.level]

        body_run = paragraph.add_run()
        body_run.text = f"{indicator.label} — {indicator.description}"
        body_run.font.size = Pt(12)
        body_run.font.color.rgb = INK
        body_run.font.name = "Calibri"
        paragraph.space_after = Pt(8)


def build_pptx(
    *,
    result: AnalyticsResult,
    charts: dict[str, bytes],
    insights: InsightBundle | None,
    organization_name: str,
    title: str | None = None,
    include_charts: bool = True,
) -> bytes:
    """Render the executive PowerPoint deck."""
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT

    deck_title = title or "Executive Business Review"
    _title_slide(presentation, result, organization_name, deck_title)
    _kpi_slide(presentation, result)

    if include_charts:
        _chart_slide(presentation, "Revenue and profit trend", [charts["revenue_trend"]])
        _chart_slide(presentation, "Performance by region and category",
                     [charts["region"], charts["category"]])
        _chart_slide(presentation, "Margin and customer mix",
                     [charts["margin_category"], charts["segments"]])
        _chart_slide(presentation, "Most profitable products", [charts["top_products"]])

    _health_slide(presentation, result)

    if result.anomalies.anomalies:
        _bullet_slide(
            presentation,
            "Detected anomalies",
            [
                f"{anomaly.period}: {anomaly.description}"
                for anomaly in result.anomalies.anomalies[:6]
            ],
            f"{result.anomalies.points_analysed} "
            f"{result.anomalies.granularity.value} periods analysed",
        )

    if insights is not None:
        note = (
            "Template-generated placeholder narrative — no AI provider configured"
            if insights.degraded
            else f"Written by {insights.model} from aggregated metrics only"
        )
        for heading, insight in (
            ("Executive summary", insights.executive_summary),
            ("Root-cause analysis", insights.root_cause),
            ("Strategic recommendations", insights.recommendations),
            ("Risks to watch", insights.risks),
        ):
            if insight is None:
                continue
            _bullet_slide(presentation, heading, plain_bullets(insight.content), note)

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer.read()


__all__ = ["build_pptx"]
